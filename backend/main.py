"""
backend/main.py
FastAPI application: all 10 REST endpoints + SSE streaming.
"""
from __future__ import annotations

import asyncio
import csv
import io
import json
import logging
import uuid
from contextlib import asynccontextmanager
from datetime import datetime
from typing import Any, Dict, List, Optional

import aiofiles
from fastapi import BackgroundTasks, FastAPI, File, Form, HTTPException, Query, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from sse_starlette.sse import EventSourceResponse

from backend.config import settings
from backend.integrations.snowflake_client import SnowflakeClient
from backend.integrations.twilio_client import TwilioClient
from backend.models.schemas import (
    APIEnvelope,
    CriterionWeight,
    JobStatus,
    LeaderboardEntry,
    RubricConfig,
    TeamRecord,
    TeamRegistration,
)
from backend.pipeline.crew_runner import EvaluationPipeline, JOB_REGISTRY, SSE_QUEUES, pipeline
import backend.pipeline.crew_runner as runner_module
from backend.pipeline.similarity import SimilarityEngine

logger = logging.getLogger(__name__)

# ── In-memory stores ──────────────────────────────────────────────────────────
TEAMS: Dict[str, TeamRecord] = {}          # team_id → TeamRecord
BLIND_COUNTER = {"n": 0}                   # sequential alias counter
_CONFIG: RubricConfig = RubricConfig()     # current rubric config
_DB: Optional[SnowflakeClient] = None
_TWILIO: Optional[TwilioClient] = None
_SIMILARITY: Optional[SimilarityEngine] = None

# ── Lifespan ──────────────────────────────────────────────────────────────────
@asynccontextmanager
async def lifespan(app: FastAPI):
    global _DB, _TWILIO, _SIMILARITY
    logger.info("EvalAI starting up …")
    _DB = SnowflakeClient()
    _TWILIO = TwilioClient()
    _SIMILARITY = SimilarityEngine()
    runner_module.pipeline = EvaluationPipeline(db_client=_DB, twilio_client=_TWILIO)

    # Restore teams from DB
    try:
        for row in _DB.get_all_teams():
            tr = TeamRecord(
                team_id=row["team_id"],
                team_name=row["team_name"],
                github_url=row["github_url"],
                registered_at=datetime.fromisoformat(str(row["registered_at"])),
                blind_alias=row.get("blind_alias", ""),
            )
            TEAMS[tr.team_id] = tr
            BLIND_COUNTER["n"] = max(BLIND_COUNTER["n"], int(tr.blind_alias.split()[-1]) if tr.blind_alias.startswith("Team ") else 0)
    except Exception as exc:
        logger.warning("Could not restore teams from DB: %s", exc)

    yield
    if _DB:
        _DB.close()
    logger.info("EvalAI shut down.")


# ── App factory ───────────────────────────────────────────────────────────────
app = FastAPI(
    title="EvalAI",
    description="AI-powered hackathon evaluation system",
    version="1.0.0",
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000", "*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ── Helpers ───────────────────────────────────────────────────────────────────
def ok(data: Any, meta: dict = None) -> dict:
    return {"status": "ok", "data": data, "meta": meta or {}}


def err(msg: str, code: int = 400):
    return JSONResponse(status_code=code, content={"status": "error", "data": None, "meta": {"message": msg}})


def _next_alias() -> str:
    BLIND_COUNTER["n"] += 1
    return f"Team {BLIND_COUNTER['n']:02d}"


def _extract_slide_text(content: bytes, filename: str) -> str:
    """Extract text from PDF or PPTX bytes. Falls back to UTF-8 decode."""
    fname = filename.lower()
    try:
        if fname.endswith(".pdf"):
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(content))
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        elif fname.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(parts)
        else:
            return content.decode("utf-8", errors="replace")
    except Exception as exc:
        logger.warning("Slide text extraction failed (%s): %s", filename, exc)
        return content.decode("utf-8", errors="replace")


# ── POST /api/teams ───────────────────────────────────────────────────────────
@app.post("/api/teams")
async def register_team(
    background_tasks: BackgroundTasks,
    team_name: str = Form(...),
    github_url: str = Form(...),
    file: Optional[UploadFile] = File(None),
):
    team_id = str(uuid.uuid4())
    alias = _next_alias()

    slide_text = ""
    if file and file.filename:
        content = await file.read()
        # Save to disk
        import os
        os.makedirs("data/uploads", exist_ok=True)
        save_path = f"data/uploads/{team_id}_{file.filename}"
        async with aiofiles.open(save_path, "wb") as f:
            await f.write(content)
        slide_text = _extract_slide_text(content, file.filename)

    team = TeamRecord(
        team_id=team_id,
        team_name=team_name,
        github_url=github_url,
        registered_at=datetime.utcnow(),
        blind_alias=alias,
        slide_text=slide_text,
    )
    TEAMS[team_id] = team

    # Index slide text for RAG (background)
    if slide_text:
        from backend.pipeline.rag import rag_pipeline
        background_tasks.add_task(rag_pipeline.index_team, team_id, slide_text)

    # Persist to DB (background)
    if _DB:
        background_tasks.add_task(
            _DB.save_team, team_id, team_name, github_url, team.registered_at, alias
        )
        background_tasks.add_task(
            _DB.log_audit, "team_registered", team_id, None,
            f"Team: {team_name}, GitHub: {github_url}", f"Alias: {alias}"
        )

    return ok({"team_id": team_id, "blind_alias": alias, "team_name": team_name})


@app.delete("/api/teams/{team_id}")
async def delete_team(team_id: str, background_tasks: BackgroundTasks):
    team = TEAMS.pop(team_id, None)
    if not team:
        return err(f"Team {team_id} not found", 404)
    if _DB:
        background_tasks.add_task(_DB.delete_team, team_id)
        background_tasks.add_task(
            _DB.log_audit, "team_deleted", team_id, None,
            f"Team {team.team_name} deleted", None
        )
        
    # Recalculate BLIND_COUNTER
    if not TEAMS:
        BLIND_COUNTER["n"] = 0
    else:
        max_n = 0
        for t in TEAMS.values():
            if t.blind_alias.startswith("Team "):
                try:
                    max_n = max(max_n, int(t.blind_alias.split()[-1]))
                except ValueError:
                    pass
        BLIND_COUNTER["n"] = max_n

    return ok({"message": "Team deleted"})


# ── POST /api/evaluate/{team_id} ──────────────────────────────────────────────
@app.post("/api/evaluate/{team_id}")
async def evaluate_team(team_id: str, background_tasks: BackgroundTasks):
    team = TEAMS.get(team_id)
    if not team:
        return err(f"Team {team_id} not found", 404)

    pipe = runner_module.pipeline
    if not pipe:
        return err("Pipeline not initialized", 500)

    job_id = pipe.create_job(team_id)
    background_tasks.add_task(_run_pipeline, job_id, team, _CONFIG)

    return ok({"job_id": job_id, "team_id": team_id, "status": "running"})


async def _run_pipeline(job_id: str, team: TeamRecord, config: RubricConfig):
    pipe = runner_module.pipeline
    if pipe:
        await pipe.run(job_id, team, config)


# ── POST /api/evaluate/batch ──────────────────────────────────────────────────
@app.post("/api/evaluate/batch")
async def evaluate_batch(background_tasks: BackgroundTasks):
    if not TEAMS:
        return err("No teams registered", 400)

    pipe = runner_module.pipeline
    if not pipe:
        return err("Pipeline not initialized", 500)

    jobs = []
    for team in TEAMS.values():
        job_id = pipe.create_job(team.team_id)
        background_tasks.add_task(_run_pipeline, job_id, team, _CONFIG)
        jobs.append({"job_id": job_id, "team_id": team.team_id})

    return ok({"batch_jobs": jobs, "total": len(jobs)})


# ── GET /api/evaluate/{job_id}/status  (SSE) ──────────────────────────────────
@app.get("/api/evaluate/{job_id}/status")
async def evaluation_status(job_id: str):
    job = JOB_REGISTRY.get(job_id)
    if not job:
        return err(f"Job {job_id} not found", 404)

    async def event_generator():
        q = SSE_QUEUES.get(job_id)

        # If job already complete, yield final status immediately
        if job.status in (JobStatus.COMPLETE, JobStatus.ERROR):
            yield {
                "event": "pipeline_complete",
                "data": json.dumps({
                    "job_id": job_id,
                    "status": job.status.value,
                    "final_score": job.scorecard.chief_judge.final_score if job.scorecard and job.scorecard.chief_judge else None,
                }),
            }
            return

        if not q:
            yield {"event": "error", "data": json.dumps({"message": "No SSE queue for job"})}
            return

        while True:
            try:
                event = await asyncio.wait_for(q.get(), timeout=30.0)
                yield {
                    "event": event["event"],
                    "data": json.dumps(event["data"]),
                }
                if event["event"] == "stream_end":
                    break
            except asyncio.TimeoutError:
                yield {"event": "heartbeat", "data": json.dumps({"job_id": job_id, "status": job.status.value})}
                if job.status in (JobStatus.COMPLETE, JobStatus.ERROR):
                    break

    return EventSourceResponse(event_generator())


# ── GET /api/results/{team_id} ────────────────────────────────────────────────
@app.get("/api/results/{team_id}")
async def get_team_results(team_id: str):
    team = TEAMS.get(team_id)
    if not team:
        return err(f"Team {team_id} not found", 404)

    # Find latest completed job for this team
    team_jobs = [j for j in JOB_REGISTRY.values() if j.team_id == team_id and j.scorecard]
    if not team_jobs:
        return ok({"team_id": team_id, "status": "no_evaluation", "scorecard": None})

    latest = max(team_jobs, key=lambda j: j.completed_at or datetime.min)
    return ok(latest.scorecard.model_dump())


# ── GET /api/results ──────────────────────────────────────────────────────────
@app.get("/api/results")
async def get_all_results():
    leaderboard: List[dict] = []
    show_real = not _CONFIG.blind_mode

    for idx, (team_id, team) in enumerate(
        sorted(TEAMS.items(), key=lambda kv: kv[1].registered_at)
    ):
        team_jobs = [j for j in JOB_REGISTRY.values() if j.team_id == team_id and j.scorecard]
        entry: Dict[str, Any] = {
            "rank": idx + 1,
            "team_id": team_id,
            "team_name": team.team_name if show_real else team.blind_alias,
            "blind_alias": team.blind_alias,
            "final_score": None,
            "recommendation": None,
            "disputed_count": 0,
            "job_status": "pending",
        }

        if team_jobs:
            latest = max(team_jobs, key=lambda j: j.completed_at or datetime.min)
            sc = latest.scorecard
            if sc and sc.chief_judge:
                entry["final_score"] = sc.chief_judge.final_score
                entry["recommendation"] = sc.chief_judge.recommendation.value if sc.chief_judge.recommendation else None
                entry["disputed_count"] = len(sc.chief_judge.disputed_criteria or [])
            entry["job_status"] = latest.status.value

        leaderboard.append(entry)

    # Sort by final_score descending
    leaderboard.sort(key=lambda e: (e["final_score"] is None, -(e["final_score"] or 0)))
    for i, e in enumerate(leaderboard):
        e["rank"] = i + 1

    return ok({"leaderboard": leaderboard, "blind_mode": _CONFIG.blind_mode, "total": len(leaderboard)})


# ── GET /api/similarity ───────────────────────────────────────────────────────
@app.get("/api/similarity")
async def get_similarity(background_tasks: BackgroundTasks):
    if not TEAMS:
        return ok({"alerts": [], "threshold": _CONFIG.similarity_threshold})

    if not _CONFIG.similarity_detection:
        return ok({"alerts": [], "disabled": True})

    teams_list = [{"team_id": t.team_id, "team_name": t.team_name} for t in TEAMS.values()]
    alerts = _SIMILARITY.cross_similarity_check(teams_list, threshold=_CONFIG.similarity_threshold)

    # Persist alerts in background
    if _DB and alerts:
        for alert in alerts:
            background_tasks.add_task(_DB.save_similarity_alert, alert)

    return ok({
        "alerts": [a.model_dump() for a in alerts],
        "threshold": _CONFIG.similarity_threshold,
        "flagged_count": len([a for a in alerts if a.flagged]),
    })


# ── GET /api/audit ────────────────────────────────────────────────────────────
@app.get("/api/audit")
async def get_audit(
    team_id: Optional[str] = Query(None),
    agent_name: Optional[str] = Query(None),
    event_type: Optional[str] = Query(None),
    export_csv: bool = Query(False, alias="csv"),
):
    if not _DB:
        return ok({"entries": []})

    entries = _DB.get_audit_log(team_id=team_id, agent_name=agent_name, event_type=event_type)

    if export_csv:
        output = io.StringIO()
        writer = csv.DictWriter(output, fieldnames=["log_id", "event_type", "team_id", "agent_name", "raw_prompt", "raw_response", "created_at"])
        writer.writeheader()
        writer.writerows(entries)
        return StreamingResponse(
            io.BytesIO(output.getvalue().encode()),
            media_type="text/csv",
            headers={"Content-Disposition": "attachment; filename=evalai_audit.csv"},
        )

    return ok({"entries": entries, "total": len(entries)})


# ── PUT /api/config ───────────────────────────────────────────────────────────
@app.put("/api/config")
async def save_config(config: RubricConfig):
    global _CONFIG
    if abs(config.weights.total() - 100.0) > 0.5:
        return err(f"Rubric weights must sum to 100 (current: {config.weights.total():.1f})")
    _CONFIG = config
    if _DB:
        _DB.log_audit("config_updated", None, None, None, config.model_dump_json())
    return ok({"message": "Config saved", "config": config.model_dump()})


# ── GET /api/config ───────────────────────────────────────────────────────────
@app.get("/api/config")
async def get_config():
    return ok(_CONFIG.model_dump())


# ── Health check ──────────────────────────────────────────────────────────────
@app.get("/health")
async def health():
    return ok({
        "status": "healthy",
        "snowflake": _DB._use_snowflake if _DB else False,
        "twilio": _TWILIO._available if _TWILIO else False,
        "groq": settings.groq_available,
        "teams_loaded": len(TEAMS),
    })
